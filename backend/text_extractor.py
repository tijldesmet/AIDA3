import os
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.docx import partition_docx
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract

# OCR talen (installeer via Tesseract)
OCR_LANGS = 'eng+fra+deu+spa+ita+nld'

def extract_text(file_path):
    text = ""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        elements = partition_pdf(filename=file_path)
        text = "\n".join(e.text for e in elements if e.text)
    elif ext == ".docx":
        elements = partition_docx(filename=file_path)
        text = "\n".join(e.text for e in elements if e.text)
    elif ext == ".rtf":
        with open(file_path, 'r', errors='ignore') as f:
            raw = f.read()
        text = rtf_to_text(raw)
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        text = "\n".join(slides)
    elif ext in [".xls", ".xlsx"]:
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            sheets = []
            for name, df in df_dict.items():
                sheets.append(f"Sheet: {name}")
                sheets.append(df.to_csv(index=False))
            text = "\n".join(sheets)
        except Exception:
            text = ""
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang=OCR_LANGS)
    else:
        text = ""
    return text
